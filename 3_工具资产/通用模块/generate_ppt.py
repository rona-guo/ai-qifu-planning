"""
Generate a professional PPT master draft for AI landing solution sharing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_BLUE = RGBColor(0xDA, 0xE3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    """Add background fill to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, font='微软雅黑'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=16, color=DARK_GRAY, font='微软雅黑', line_spacing=1.5):
    """Add multi-line text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.space_after = Pt(size * 0.6)
    return txBox

def add_table(slide, left, top, width, height, data, col_widths=None, header_color=BLUE, font_size=12):
    """Add a styled table."""
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    # Style header
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = data[0][j]
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
    
    # Style data rows
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j]
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_GRAY
                p.font.name = '微软雅黑'
                p.alignment = PP_ALIGN.LEFT
    
    return table

def add_page_number(slide, num, total):
    """Add page number at bottom right."""
    add_text(slide, Inches(12.3), Inches(7.0), Inches(1.0), Inches(0.4),
             f'{num}/{total}', size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def add_title_bar(slide, title, subtitle=None):
    """Add a title bar at the top of a content slide."""
    add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.05), ORANGE)
    add_text(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
             title, size=28, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.7), Inches(10), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xAA, 0xCC, 0xEE))

TOTAL_SLIDES = 18

# ==================== Slide 1: Title ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)
# Accent bar
add_rect(slide, 0, Inches(5.5), W, Inches(0.06), ORANGE)
# Title
add_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(1.2),
         '标准企业AI落地解决方案', size=44, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(2.8), Inches(11), Inches(0.8),
         '—— 不知道从哪开始？先看这份 ——', size=22, color=RGBColor(0xAA, 0xCC, 0xEE))
# Subtitle
add_text(slide, Inches(1.0), Inches(3.8), Inches(11), Inches(0.6),
         '低门槛入场 · 高价值运维 · 持续进化', size=18, color=RGBColor(0x99, 0xBB, 0xDD))
# Footer
add_text(slide, Inches(1.0), Inches(6.0), Inches(11), Inches(0.5),
         '可整本使用  ·  可拆分组件使用  ·  可适配行业定制', size=14, color=RGBColor(0x88, 0xAA, 0xCC))

# ==================== Slide 2: 企业正在"出血" ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '你的企业可能正在"出血"')
add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
         '老板都在问：AI该用什么、怎么用、值不值得做？但先看看你的企业有没有这些问题：',
         size=16, color=GRAY)

cards = [
    ('信息流失', '客户资料散在微信、Excel里\n人走数据走', RED),
    ('流程断裂', '一份合同走5个签字\n全靠纸质流转', ORANGE),
    ('经验断层', '老员工经验在脑子里\n新人从头学起', RGBColor(0x8B, 0x69, 0x14)),
    ('决策盲区', '老板拍脑袋决策\n缺乏数据支撑', GRAY),
]
card_w = Inches(2.9)
card_h = Inches(2.2)
start_x = Inches(0.6)
gap = Inches(0.25)
for i, (title, desc, color) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    add_rounded_rect(slide, x, Inches(2.2), card_w, card_h, LIGHT_GRAY)
    add_rect(slide, x, Inches(2.2), card_w, Inches(0.08), color)
    add_text(slide, x + Inches(0.3), Inches(2.5), card_w - Inches(0.6), Inches(0.5),
             title, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.1), card_w - Inches(0.6), Inches(1.2),
             desc, size=15, color=DARK_GRAY)

add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.6),
         'AI不是万能药，但它是目前成本最低的止血工具——前提是用对方法。',
         size=18, color=DARK_BLUE, bold=True)
add_page_number(slide, 2, TOTAL_SLIDES)

# ==================== Slide 3: 三断层 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '为什么"自己搭Agent"会失败？', '三断层理论')

faults = [
    ('断层一：数据主权', '员工把公司合同、客户名单上传到公网平台\n数据泄露是不可逆的合规风险', RED),
    ('断层二：业务上下文', '通用AI不懂你的审批流程、客户等级\n生成的是"正确的废话"', ORANGE),
    ('断层三：持续运营', '员工离职，他搭的Agent就没了\n每次人员变动都是一次归零', GRAY),
]
for i, (title, desc, color) in enumerate(faults):
    y = Inches(1.6) + i * Inches(1.6)
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1.3), LIGHT_GRAY)
    add_rect(slide, Inches(0.6), y, Inches(0.1), Inches(1.3), color)
    add_text(slide, Inches(1.0), y + Inches(0.15), Inches(4), Inches(0.5),
             title, size=18, color=color, bold=True)
    add_text(slide, Inches(1.0), y + Inches(0.6), Inches(11), Inches(0.7),
             desc, size=14, color=DARK_GRAY)

add_text(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
         '结论：AI工具是"锄头"，我们提供的是"现代化农场"的整体规划与运营方案。',
         size=15, color=DARK_BLUE, bold=True)
add_page_number(slide, 3, TOTAL_SLIDES)

# ==================== Slide 4: 四步落地法 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '我们的做法：四步落地法', '不是卖你一套系统，是帮你把AI变成生产力')

steps = [
    ('第一步', 'AI诊断', '看清楚棋盘\n1-2周', BLUE),
    ('第二步', '方案设计', '画好路线图\n1周', RGBColor(0x2E, 0x74, 0xB5)),
    ('第三步', '系统实施', '把系统搭起来\n4-16周', RGBColor(0x1F, 0x4E, 0x79)),
    ('第四步', '持续运维', '越用越值钱\n持续', DARK_BLUE),
]
step_w = Inches(2.8)
step_h = Inches(2.5)
arrow_w = Inches(0.35)
start_x = Inches(0.5)
y = Inches(2.0)

for i, (step, title, desc, color) in enumerate(steps):
    x = start_x + i * (step_w + arrow_w)
    # Card
    add_rounded_rect(slide, x, y, step_w, step_h, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), step_w - Inches(0.6), Inches(0.4),
             step, size=14, color=RGBColor(0xCC, 0xDD, 0xFF))
    add_text(slide, x + Inches(0.3), y + Inches(0.7), step_w - Inches(0.6), Inches(0.6),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.4), step_w - Inches(0.6), Inches(1.0),
             desc, size=14, color=RGBColor(0xDD, 0xEE, 0xFF))
    # Arrow
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + step_w + Inches(0.02), y + Inches(0.9),
            arrow_w - Inches(0.04), Inches(0.6)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
         '你可以随时停在第一步——诊断完觉得不值得做，我们尊重你的决定。',
         size=15, color=GRAY)
add_text(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.6),
         '诊断费2-3万，如后续签约实施，诊断费可抵扣实施费。',
         size=15, color=ORANGE, bold=True)
add_page_number(slide, 4, TOTAL_SLIDES)

# ==================== Slide 5: 能力组件库 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '能力组件库：按需组合，按需定制', '已有组件+定制开发，不限于此')

modules = [
    ('A', '智能文档工厂', '报告自动生成\n合同/文档审核\n模板管理', '周报从2h→15min'),
    ('B', '企业知识中枢', '制度问答\n新人陪练\n经验沉淀\n法规查询', '新人上手3月→1月'),
    ('C', '智能运营助手', '数据汇总\n风险预警\n分析报告', '老板实时看数据'),
    ('D', '统一门户与连接', '企微/钉钉集成\n待办聚合\nAI客服入口', '零学习成本'),
]
mod_w = Inches(2.9)
mod_h = Inches(3.2)
start_x = Inches(0.5)
gap = Inches(0.25)
y = Inches(1.5)
colors = [BLUE, RGBColor(0x2E, 0x74, 0xB5), RGBColor(0x1F, 0x4E, 0x79), DARK_BLUE]

for i, (label, title, features, effect) in enumerate(modules):
    x = start_x + i * (mod_w + gap)
    color = colors[i]
    add_rounded_rect(slide, x, y, mod_w, mod_h, LIGHT_GRAY)
    add_rect(slide, x, y, mod_w, Inches(0.7), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.5),
             label, size=24, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.8), y + Inches(0.15), mod_w - Inches(1), Inches(0.5),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.9), mod_w - Inches(0.6), Inches(1.8),
             features, size=13, color=DARK_GRAY)
    # Effect badge
    add_rounded_rect(slide, x + Inches(0.2), y + Inches(2.5), mod_w - Inches(0.4), Inches(0.55), ORANGE)
    add_text(slide, x + Inches(0.3), y + Inches(2.55), mod_w - Inches(0.6), Inches(0.45),
             effect, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# "按需定制" hint below the cards
add_text(slide, Inches(0.5), Inches(5.0), Inches(12.5), Inches(0.8),
         '以上为已验证组件示例，不限于此  ·  你的场景不在这里？告诉我们，按需定制',
         size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(slide, 5, TOTAL_SLIDES)

# ==================== Slide 6: 人机协同四步闭环 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'AI怎么工作：人机协同四步闭环', 'AI是参谋不是裁判')

steps_data = [
    ('1', '规则引擎预审', '必填校验、格式检查\n过滤低级错误', BLUE),
    ('2', 'AI引擎+知识库', '语义分析、法规比对\n风险标记+置信度', RGBColor(0x2E, 0x74, 0xB5)),
    ('3', '业务人员终审', '查看原文+AI标记\n确认或驳回（强制）', ORANGE),
    ('4', '系统自动学习', '人工修正更新知识库\nAI下次更准', GREEN),
]
step_w = Inches(2.8)
step_h = Inches(2.5)
start_x = Inches(0.5)
y = Inches(1.8)

for i, (num, title, desc, color) in enumerate(steps_data):
    x = start_x + i * (step_w + Inches(0.35))
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, x + Inches(1.0), y + Inches(0.1), Inches(0.8), Inches(0.6),
             num, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Card
    add_rounded_rect(slide, x, y + Inches(1.0), step_w, Inches(2.5), LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(1.2), step_w - Inches(0.4), Inches(0.5),
             title, size=16, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.8), step_w - Inches(0.4), Inches(1.5),
             desc, size=13, color=DARK_GRAY)
    # Arrow down
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            x + step_w + Inches(0.02), y + Inches(1.5),
            Inches(0.3), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

add_text(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.6),
         '关键：AI干80%的活，人做20%的关键决策。AI越用越准，因为每次人工修正都让AI变聪明。',
         size=15, color=DARK_BLUE, bold=True)
add_page_number(slide, 6, TOTAL_SLIDES)

# ==================== Slide 7: 数据安全 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '数据安全：私有化 + 四级分级', '你的数据不会出现在任何公网上')

# Left: deployment diagram
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5), LIGHT_BLUE)
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
         '你的服务器（内网）', size=18, color=DARK_BLUE, bold=True)
add_rounded_rect(slide, Inches(1.0), Inches(2.4), Inches(4.7), Inches(2.2), WHITE)
lines = [
    ('业务平台 + AI引擎 + 知识库', False),
    ('', False),
    ('数据100%在你的机房里', True),
    ('我们远程维护，碰不到你的数据', True),
    ('', False),
    ('为什么不是SaaS？', True),
    ('合同、客户名单上了公网就收不回来', False),
    ('这是不可逆的合规风险', False),
]
add_multiline(slide, Inches(1.2), Inches(2.6), Inches(4.3), Inches(2.0), lines, size=14, color=DARK_GRAY)

# Right: L1-L4 table
table_data = [
    ['级别', '什么数据', '怎么处理'],
    ['L1 公开', '宣传资料', '可公开'],
    ['L2 内部', '周报、通知', '内部流转'],
    ['L3 敏感', '合同、客户信息', '脱敏处理'],
    ['L4 核心', '财报、商业机密', '物理隔离'],
]
add_text(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.5),
         '数据分级管理：L1-L4', size=18, color=DARK_BLUE, bold=True)
add_table(slide, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.8), table_data,
          col_widths=[Inches(1.5), Inches(2.5), Inches(2.0)], font_size=13)

add_text(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.6),
         '平台选择不绑定：明道云/简道云/宜搭/定制开发——根据你的情况选最合适的。',
         size=14, color=GRAY)
add_page_number(slide, 7, TOTAL_SLIDES)

# ==================== Slide 8: 透明报价 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '投入多少：透明报价', '成本透明，不赚差价')

# Cost breakdown
table_data = [
    ['费用项', '金额', '性质'],
    ['平台许可', '6,000-20,000元/年', '年付，你的资产'],
    ['服务器硬件', '8,000-12,000元', '一次性，你的资产'],
    ['实施费（基础版）', '3-5万', '1-2个组件+培训'],
    ['实施费（标准版）', '5-8万', '3-4个组件+AI集成'],
    ['实施费（旗舰版）', '8-12万', '多组件+深度定制'],
    ['年度运维', '0.8-2万/年', '监控+AI运营'],
    ['培训费', '0.8-2.5万', '可选采购'],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(7.5), Inches(4.5), table_data,
          col_widths=[Inches(2.5), Inches(2.5), Inches(2.5)], font_size=12)

# Typical scenario
add_rounded_rect(slide, Inches(8.5), Inches(1.4), Inches(4.3), Inches(4.5), DARK_BLUE)
add_text(slide, Inches(8.7), Inches(1.6), Inches(4), Inches(0.5),
         '典型场景', size=18, color=WHITE, bold=True)
add_text(slide, Inches(8.7), Inches(2.1), Inches(4), Inches(0.4),
         '50人贸易公司', size=14, color=RGBColor(0xAA, 0xCC, 0xEE))

scenario = [
    ('平台许可（年）', '6,000元'),
    ('服务器', '10,000元'),
    ('实施费（标准版）', '60,000元'),
    ('首年运维', '12,000元'),
]
for i, (item, amount) in enumerate(scenario):
    y = Inches(2.6) + i * Inches(0.5)
    add_text(slide, Inches(8.7), y, Inches(2.5), Inches(0.4),
             item, size=13, color=RGBColor(0xBB, 0xCC, 0xDD))
    add_text(slide, Inches(11.0), y, Inches(1.5), Inches(0.4),
             amount, size=13, color=WHITE, align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(8.7), Inches(4.7), Inches(3.9), Inches(0.04), ORANGE)
add_text(slide, Inches(8.7), Inches(4.8), Inches(2), Inches(0.5),
         '首期总投入', size=14, color=RGBColor(0xBB, 0xCC, 0xDD))
add_text(slide, Inches(10.0), Inches(4.8), Inches(2.6), Inches(0.5),
         '约88,000元', size=20, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)
add_text(slide, Inches(8.7), Inches(5.4), Inches(4), Inches(0.4),
         '10万以内，跑起一套私有化AI系统', size=13, color=RGBColor(0x99, 0xBB, 0xDD))

add_page_number(slide, 8, TOTAL_SLIDES)

# ==================== Slide 9: 为什么选我们 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '为什么选我们', '六个关键差异')

comparisons = [
    ('传统软件公司', '30-50万起，3-6个月', '我们', '10万以内，4-16周'),
    ('SaaS平台', '数据在公网，你是租户', '我们', '私有化，数据是你的'),
    ('AI工具商', '卖完工具就走', '我们', '卖完才刚开始——持续运维'),
    ('培训公司', '只讲理论不落地', '我们', '先诊断再实施，有系统'),
    ('自己搭', '数据泄露不可逆', '我们', '专业实施+合规保障'),
]
for i, (left_t, left_d, right_t, right_d) in enumerate(comparisons):
    y = Inches(1.5) + i * Inches(1.05)
    # Left card
    add_rounded_rect(slide, Inches(0.6), y, Inches(5.3), Inches(0.85), RGBColor(0xFB, 0xEA, 0xEA))
    add_text(slide, Inches(0.8), y + Inches(0.08), Inches(2), Inches(0.4),
             left_t, size=14, color=RED, bold=True)
    add_text(slide, Inches(0.8), y + Inches(0.45), Inches(4.8), Inches(0.4),
             left_d, size=13, color=GRAY)
    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), y + Inches(0.25), Inches(0.5), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()
    # Right card
    add_rounded_rect(slide, Inches(6.6), y, Inches(6.2), Inches(0.85), RGBColor(0xEA, 0xF5, 0xEA))
    add_text(slide, Inches(6.8), y + Inches(0.08), Inches(1.5), Inches(0.4),
             right_t, size=14, color=GREEN, bold=True)
    add_text(slide, Inches(6.8), y + Inches(0.45), Inches(5.7), Inches(0.4),
             right_d, size=13, color=DARK_GRAY)

add_text(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
         '核心差异：我们不卖工具，我们卖"AI变成生产力"的能力。',
         size=16, color=DARK_BLUE, bold=True)
add_page_number(slide, 9, TOTAL_SLIDES)

# ==================== Slide 10: 合作流程 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '合作流程：从接触到上线')

flow = [
    ('Step 1', '初步沟通', '免费', '你说情况→我们判断\n是否适合AI'),
    ('Step 2', 'AI诊断', '2-3万\n可抵扣', '1-2周→出诊断报告\n+路线图'),
    ('Step 3', '方案设计', '免费\n含在实施费', '1周→出方案书\n+费用明细→签合同'),
    ('Step 4', '系统实施', '3-12万', '4-16周\n分阶段交付'),
    ('Step 5', '培训赋能', '0.8-\n2.5万', '2-4天\n全员培训'),
    ('Step 6', '上线运维', '年度\n签约', '正式上线\n季度复盘'),
]
step_w = Inches(1.9)
start_x = Inches(0.4)
y = Inches(2.0)
colors_flow = [GRAY, BLUE, RGBColor(0x2E, 0x74, 0xB5), RGBColor(0x1F, 0x4E, 0x79), ORANGE, GREEN]

for i, (step, title, cost, desc) in enumerate(flow):
    x = start_x + i * (step_w + Inches(0.15))
    color = colors_flow[i]
    add_rounded_rect(slide, x, y, step_w, Inches(3.5), LIGHT_GRAY)
    add_rect(slide, x, y, step_w, Inches(0.6), color)
    add_text(slide, x + Inches(0.1), y + Inches(0.05), step_w - Inches(0.2), Inches(0.5),
             step, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.7), step_w - Inches(0.2), Inches(0.4),
             title, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(1.2), step_w - Inches(0.2), Inches(0.8),
             cost, size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(2.1), step_w - Inches(0.2), Inches(1.3),
             desc, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.6),
         '你可以随时停在Step 2——诊断完觉得不值得做，诊断报告你留着。',
         size=15, color=GRAY)
add_page_number(slide, 10, TOTAL_SLIDES)

# ==================== Slide 11: 风险控制 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '我们承诺什么，不承诺什么', '诚实是最好的销售策略')

# Left: commitments
add_text(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(0.5),
         '我们承诺的', size=20, color=GREEN, bold=True)
commitments = [
    '数据安全：100%私有化部署，数据物理隔离在内网',
    '交付标准：每个组件UAT验收，客户签字确认',
    '质量保障：交付后30天免费缺陷修复',
    '成本透明：所有费用明码标价，不赚差价',
    '需求冻结：签约后冻结范围，变更走增项流程',
    '持续服务：年度运维，不交付完就走',
]
for i, text in enumerate(commitments):
    y = Inches(2.0) + i * Inches(0.55)
    add_text(slide, Inches(0.8), y, Inches(0.3), Inches(0.4), '✓', size=16, color=GREEN, bold=True)
    add_text(slide, Inches(1.1), y, Inches(5.3), Inches(0.4), text, size=13, color=DARK_GRAY)

# Right: not committed
add_text(slide, Inches(7.0), Inches(1.4), Inches(6), Inches(0.5),
         '我们不承诺的', size=20, color=RED, bold=True)
not_committed = [
    ('AI 100%准确', '任何承诺100%的都是骗子\n我们用人机协同兜底'),
    ('一周速成', '好的系统需要时间打磨\n我们不赶工出次品'),
    ('全自动无人化', 'AI是参谋不是裁判\n关键决策必须有人参与'),
]
for i, (title, desc) in enumerate(not_committed):
    y = Inches(2.0) + i * Inches(1.3)
    add_text(slide, Inches(7.2), y, Inches(0.3), Inches(0.4), '✗', size=16, color=RED, bold=True)
    add_text(slide, Inches(7.5), y, Inches(5), Inches(0.4),
             title, size=14, color=DARK_GRAY, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.4), Inches(5), Inches(0.7),
             desc, size=12, color=GRAY)

add_page_number(slide, 11, TOTAL_SLIDES)

# ==================== Slide 12: 方案延伸 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '方案延伸：从一个标准方案到行业专属', '母体方案 + 行业适配 = 行业专属方案')

# Center: standard solution
add_rounded_rect(slide, Inches(4.7), Inches(1.8), Inches(4), Inches(1.2), DARK_BLUE)
add_text(slide, Inches(4.9), Inches(1.95), Inches(3.6), Inches(0.5),
         '标准AI落地方案', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.9), Inches(2.45), Inches(3.6), Inches(0.4),
         '组件库 + 技术架构 + 运维体系', size=13, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# Extensions
extensions = [
    ('人力资源服务', '已有初版', GREEN),
    ('劳动关系风控', '参赛中', BLUE),
    ('制造业', '可定制', GRAY),
    ('建工行业', '可定制', GRAY),
]
for i, (name, status, color) in enumerate(extensions):
    y = Inches(3.5) + i * Inches(0.8)
    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.0), y + Inches(0.1), Inches(0.6), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()
    # Left: industry
    add_rounded_rect(slide, Inches(1.5), y, Inches(2.5), Inches(0.6), LIGHT_GRAY)
    add_text(slide, Inches(1.7), y + Inches(0.1), Inches(2.2), Inches(0.4),
             name, size=14, color=DARK_GRAY, bold=True)
    # Right: result
    add_rounded_rect(slide, Inches(8.7), y, Inches(3.5), Inches(0.6), color)
    add_text(slide, Inches(8.9), y + Inches(0.1), Inches(3.2), Inches(0.4),
             status, size=13, color=WHITE, bold=True)

add_text(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
         '底层架构不变，上层组件按行业适配。你的行业不在列表里？告诉我们，我们可以定制。',
         size=14, color=GRAY)
add_page_number(slide, 12, TOTAL_SLIDES)

# ==================== Slide 13: 知识沉淀 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '交付后：系统越用越值钱', '传统交付"做完就结束"，我们做的是"做完才开始"')

# Flow
flow_items = [
    ('项目交付', BLUE),
    ('AI使用数据', RGBColor(0x2E, 0x74, 0xB5)),
    ('知识库更新', RGBColor(0x1F, 0x4E, 0x79)),
    ('Prompt调优', ORANGE),
    ('效果提升', GREEN),
    ('价值放大', DARK_BLUE),
]
for i, (text, color) in enumerate(flow_items):
    x = Inches(0.5) + i * Inches(2.1)
    add_rounded_rect(slide, x, Inches(1.8), Inches(1.9), Inches(0.8), color)
    add_text(slide, x + Inches(0.1), Inches(1.95), Inches(1.7), Inches(0.5),
             text, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.92), Inches(2.0), Inches(0.15), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# Deposit content table
table_data = [
    ['沉淀内容', '来源', '用途'],
    ['业务模板库', '客户业务文档', '下个同行业客户直接复用'],
    ['Prompt优化记录', '使用过程调优', '避免重复踩坑'],
    ['行业知识库', '客户授权的法规/制度', '同行业客户共享'],
    ['交付SOP', '实施过程记录', '标准化流程，缩短周期'],
    ['效果数据', 'AI使用统计', '向新客户展示量化效果'],
]
add_table(slide, Inches(0.6), Inches(3.2), Inches(12.2), Inches(3.0), table_data,
          col_widths=[Inches(2.5), Inches(4.0), Inches(5.7)], font_size=12)

add_page_number(slide, 13, TOTAL_SLIDES)

# ==================== Slide 14: 常见问题 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '常见问题')

faqs = [
    ('连OA都没有，能用AI吗？', '恰恰最需要。我们先帮你理顺数据、跑通流程（建底座），再叠加AI。不是空中楼阁。'),
    ('预算紧，能不能先试试？', '可以。诊断2-3万先看值不值得做。最低3万跑起一个组件。平台年费6,000，服务器1万以内。'),
    ('AI会不会取代我们的人？', '不会。AI干机械重复的活，人做关键决策。让你的团队从体力活中解放出来。'),
    ('自己用Coze搭不行吗？', '个人用没问题。但企业级要解决三件事：数据安全、业务贴合、持续运营。我们解决的就是这三件事。'),
    ('系统上线后谁管？', '年度运维。不用招IT，不用担心没人管。知识库我们更新，Prompt我们调优。'),
]
for i, (q, a) in enumerate(faqs):
    y = Inches(1.4) + i * Inches(1.1)
    add_text(slide, Inches(0.6), y, Inches(0.4), Inches(0.4),
             'Q', size=16, color=BLUE, bold=True)
    add_text(slide, Inches(1.0), y, Inches(11.5), Inches(0.4),
             q, size=15, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(0.6), y + Inches(0.45), Inches(0.4), Inches(0.4),
             'A', size=16, color=ORANGE, bold=True)
    add_text(slide, Inches(1.0), y + Inches(0.45), Inches(11.5), Inches(0.5),
             a, size=13, color=DARK_GRAY)

add_page_number(slide, 14, TOTAL_SLIDES)

# ==================== Slide 15: 关于我们 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '关于我们')

add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
         '专注于中小企业AI落地的科技服务公司', size=20, color=DARK_BLUE, bold=True)

features = [
    ('产品经理出身', '用业务语言不是技术术语'),
    ('BD+咨询+数字化', '复合背景，懂老板要什么'),
    ('AI工具自用', '一个人做到传统团队的深度'),
    ('专家网络', '按项目协作，不养冗余团队'),
]
for i, (title, desc) in enumerate(features):
    y = Inches(2.5) + i * Inches(0.8)
    add_text(slide, Inches(1.0), y, Inches(0.3), Inches(0.4), '▸', size=16, color=ORANGE, bold=True)
    add_text(slide, Inches(1.4), y, Inches(4), Inches(0.4),
             title, size=16, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(5.5), y, Inches(7), Inches(0.4),
             desc, size=14, color=GRAY)

add_text(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.6),
         '核心理念：不是卖工具，是帮老板把AI变成生产力。',
         size=18, color=DARK_BLUE, bold=True)
add_page_number(slide, 15, TOTAL_SLIDES)

# ==================== Slide 16: 一页纸总览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '一页纸总览', '一页看懂全部')

# 4 columns
col_w = Inches(3.0)
gap = Inches(0.15)
start_x = Inches(0.5)

col_data = [
    ('做什么', BLUE, [
        'AI诊断：看清楚棋盘',
        '方案设计：画好路线图',
        '系统实施：把系统搭起来',
        '持续运维：越用越值钱',
    ]),
    ('卖什么', RGBColor(0x2E, 0x74, 0xB5), [
        'A 智能文档工厂',
        'B 企业知识中枢',
        'C 智能运营助手',
        'D 统一门户与连接',
    ]),
    ('凭什么', RGBColor(0x1F, 0x4E, 0x79), [
        '私有化部署，数据安全',
        '人机协同，AI干活人把关',
        'L1-L4数据分级管理',
        '年度运维，持续进化',
    ]),
    ('多少钱', DARK_BLUE, [
        '诊断：2-3万',
        '实施：3-12万',
        '运维：0.8-2万/年',
        '培训：0.8-2.5万',
        '',
        '典型：10万以内跑起来',
    ]),
]
for i, (title, color, items) in enumerate(col_data):
    x = start_x + i * (col_w + gap)
    add_rounded_rect(slide, x, Inches(1.4), col_w, Inches(5.2), LIGHT_GRAY)
    add_rect(slide, x, Inches(1.4), col_w, Inches(0.6), color)
    add_text(slide, x + Inches(0.1), Inches(1.45), col_w - Inches(0.2), Inches(0.5),
             title, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        if item:
            y = Inches(2.2) + j * Inches(0.55)
            add_text(slide, x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.4),
                     '• ' + item, size=13, color=DARK_GRAY)

add_page_number(slide, 16, TOTAL_SLIDES)

# ==================== Slide 17: 下一步行动 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, 0, Inches(5.5), W, Inches(0.06), ORANGE)

add_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.8),
         '下一步：从诊断开始', size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(0.6),
         '不确定值不值得做？先花2-3万做诊断，看清楚再决定。',
         size=20, color=RGBColor(0xAA, 0xCC, 0xEE))

steps_text = [
    ('1', '联系我们，说你的情况'),
    ('2', '我们判断是否适合AI，给你建议'),
    ('3', '签订诊断协议，1-2周出报告'),
    ('4', '你看了报告，决定做不做'),
]
for i, (num, text) in enumerate(steps_text):
    y = Inches(3.3) + i * Inches(0.55)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text(slide, Inches(1.0), y + Inches(0.02), Inches(0.4), Inches(0.35),
             num, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.6), y + Inches(0.02), Inches(10), Inches(0.4),
             text, size=16, color=WHITE)

add_text(slide, Inches(1.0), Inches(6.0), Inches(11), Inches(0.5),
         '诊断费可抵扣实施费。不签约也不亏——诊断报告你留着。',
         size=15, color=ORANGE, bold=True)

# ==================== Slide 18: Thank You ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, 0, Inches(3.5), W, Inches(0.06), ORANGE)

add_text(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(1.0),
         'AI不是万能药\n但它是目前成本最低的止血工具', size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(0.6),
         '—— 前提是用对方法 ——', size=18, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.5),
         '联系我们  ·  开始你的AI落地之旅', size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = r'D:\AI落地\4_业务运营\展示物料\AI落地标准方案分享PPT母稿.pptx'
prs.save(output_path)
print(f'PPT saved: {output_path}')
print(f'Total slides: {len(prs.slides)}')
